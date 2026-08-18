"""
이미지로만 된 슬라이드를 읽을 때, **동시에 부르되 순서는 지키는가**에 대한 테스트.

배경(2026-08-16 실측): 글자가 없는 10장짜리 PPT를 업로드했더니 **추출에만 116초**가 걸렸다.
정작 대본 생성은 12초다 — 사용자가 기다리는 시간의 대부분이 이미지를 한 장씩 순서대로
읽는 데 쓰이고 있었다. 이미지들끼리는 서로 의존하지 않으므로 한꺼번에 보내면 된다.

병렬로 바꾸면서 생기는 위험이 정확히 하나다. **순서가 뒤섞이는 것.**
슬라이드 본문은 "텍스트박스 → 큰 이미지 → 작은 이미지" 순으로 이어 붙는데, 완료되는
순서대로 담으면 문장이 뒤죽박죽이 된다. 대본 생성은 이 텍스트를 근거로 삼으므로
조용히 이상한 대본이 나올 뿐 에러는 안 난다 — 그래서 여기서 못 박는다.
"""
import io
import threading
import time

import pytest
from pptx import Presentation
from pptx.util import Inches

from utils.ppt_extractor import PptExtractor


def _png(width, height, color=(120, 140, 200)):
    """python-pptx가 받아줄 최소 PNG 바이트."""
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (width, height), color).save(buf, format="PNG")
    buf.seek(0)
    return buf


def _deck(tmp_path, slides):
    """slides = [(텍스트 or None, [(가로, 세로), ...])] → .pptx 파일 경로."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for text, images in slides:
        slide = prs.slides.add_slide(blank)
        if text:
            box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(6), Inches(1))
            box.text_frame.text = text
        for index, (w, h) in enumerate(images):
            slide.shapes.add_picture(
                _png(w, h), Inches(0.5), Inches(1.5 + index * 1.2), Inches(3), Inches(1)
            )
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return str(path)


class _Recorder:
    """비전 호출을 가로채, 이미지 크기별로 정해진 문자열을 돌려준다."""

    def __init__(self, delay=0.0, fail_on=None):
        self.delay = delay
        self.fail_on = fail_on
        self.calls = []
        self.concurrent = 0
        self.peak = 0
        self._lock = threading.Lock()

    def extract_text_from_image(self, blob, context_hint, content_type):
        from PIL import Image

        size = Image.open(io.BytesIO(blob)).size
        with self._lock:
            self.calls.append(size)
            self.concurrent += 1
            self.peak = max(self.peak, self.concurrent)
        try:
            if self.delay:
                time.sleep(self.delay)
            if self.fail_on and size == self.fail_on:
                raise RuntimeError("비전 호출 실패")
            return "이미지%dx%d" % size
        finally:
            with self._lock:
                self.concurrent -= 1


def _extract(tmp_path, slides, recorder):
    extractor = PptExtractor()
    extractor.image_text_extractor = recorder
    return extractor.extract_structured_data(_deck(tmp_path, slides))


# ── 순서 ────────────────────────────────────────────────────────────────────

def test_textbox_comes_before_image_text(tmp_path):
    """텍스트박스 내용이 이미지에서 읽은 것보다 앞에 와야 한다."""
    result = _extract(tmp_path, [("제목", [(400, 300)])], _Recorder())
    content = result["slides"][0]["content"]

    assert content.index("제목") < content.index("이미지400x300")


def test_bigger_image_comes_first(tmp_path):
    """넓이 순으로 읽는 규칙은 병렬로 바꿔도 유지돼야 한다(큰 이미지가 본문일 확률이 높다)."""
    slides = [(None, [(200, 200), (600, 500), (300, 300)])]
    result = _extract(tmp_path, slides, _Recorder())
    lines = result["slides"][0]["content"].splitlines()

    assert lines == ["이미지600x500", "이미지300x300", "이미지200x200"]


def test_slides_keep_their_own_text(tmp_path):
    """완료 순서가 뒤섞여도 각 슬라이드는 자기 이미지 결과만 가져야 한다.

    앞 슬라이드가 느리고 뒤 슬라이드가 빠르면 결과가 역순으로 도착한다 —
    그때 내용이 남의 슬라이드로 새면 대본이 통째로 어긋난다.
    """
    slides = [(None, [(500, 400)]), (None, [(450, 400)]), (None, [(400, 400)])]
    result = _extract(tmp_path, slides, _Recorder(delay=0.05))

    assert [s["content"] for s in result["slides"]] == [
        "이미지500x400", "이미지450x400", "이미지400x400",
    ]
    assert [s["slide_number"] for s in result["slides"]] == [1, 2, 3]


# ── 병렬성 ──────────────────────────────────────────────────────────────────

def test_images_are_read_concurrently(tmp_path):
    """실제로 동시에 나가야 한다. 순차로 돌아가면 이 테스트가 실패한다."""
    slides = [(None, [(400 + i, 400)]) for i in range(6)]
    recorder = _Recorder(delay=0.15)

    started = time.time()
    _extract(tmp_path, slides, recorder)
    elapsed = time.time() - started

    assert recorder.peak > 1, "동시에 실행된 호출이 없습니다 — 순차 처리로 되돌아갔습니다"
    # 순차면 6 × 0.15 = 0.9초 이상. 동시 4개면 두 번에 나눠 0.3초대.
    assert elapsed < 0.75, "순차 실행에 가까운 시간이 걸렸습니다 (%.2f초)" % elapsed


def test_all_images_are_read_exactly_once(tmp_path):
    """묶어서 보내다 빠뜨리거나 두 번 보내면 안 된다(유료 호출이다)."""
    slides = [(None, [(500, 400), (450, 400)]), ("글자만 있는 장", []), (None, [(600, 400)])]
    recorder = _Recorder()
    _extract(tmp_path, slides, recorder)

    assert sorted(recorder.calls) == [(450, 400), (500, 400), (600, 400)]


# ── 안전장치 ────────────────────────────────────────────────────────────────

def test_one_failed_image_does_not_break_the_upload(tmp_path):
    """이미지 한 장이 실패해도 나머지는 살아야 한다. 업로드 전체가 죽으면 안 된다."""
    slides = [(None, [(600, 400)]), (None, [(500, 400)])]
    result = _extract(tmp_path, slides, _Recorder(fail_on=(600, 400)))

    # 실패한 장도 **페이지는 남는다** — 빈 내용일 뿐 목록에서 빠지면 안 된다.
    # (예전엔 여기서 빠지는 걸 정상으로 봤는데, 그 규칙이 14장 PPT를 10장으로 만들었다. 2026-08-18)
    contents = {s["slide_number"]: s["content"] for s in result["slides"]}
    assert contents.get(2) == "이미지500x400"
    assert contents.get(1) == ""


def test_text_heavy_slides_skip_vision_entirely(tmp_path):
    """글자가 충분한 슬라이드는 이미지를 읽지 않는다 — 유료 호출을 아끼는 규칙이다."""
    long_text = "이 슬라이드에는 충분히 긴 설명이 이미 텍스트로 들어 있습니다. " * 2
    recorder = _Recorder()
    result = _extract(tmp_path, [(long_text, [(600, 400)])], recorder)

    assert recorder.calls == []
    assert "이미지" not in result["slides"][0]["content"]


def test_deck_without_images_needs_no_vision(tmp_path):
    """이미지가 없으면 스레드풀을 아예 열지 않는다."""
    recorder = _Recorder()
    result = _extract(tmp_path, [("첫 장입니다.", []), ("둘째 장입니다.", [])], recorder)

    assert recorder.calls == []
    assert [s["content"] for s in result["slides"]] == ["첫 장입니다.", "둘째 장입니다."]


# ---------------------------------------------------------------------------
# 내용이 안 뽑힌 장도 페이지는 유지해야 한다 (2026-08-18)
#
# 실측: 02_이미지형.pptx 14장을 올렸더니 10장만 등록됐다. 텍스트도 없고 이미지도
# 걸러진(장식 판정/비전 실패) 장을 추출기가 통째로 버렸기 때문이다. 사용자에게는
# "내 PPT에서 장이 사라졌다"로 보이고, 장 번호가 밀리면 썸네일·하이라이트 좌표까지
# 연쇄로 어긋난다. 빈 장의 대본은 생성기의 근거 없는 장 처리(_is_thin_source)가 맡는다.
# ---------------------------------------------------------------------------


def test_empty_slides_are_kept_with_original_numbers(tmp_path, monkeypatch):
    """텍스트도 이미지도 없는 장이 사라지면 안 되고, 번호도 원본과 같아야 한다."""
    path = _deck(tmp_path, [
        ("첫 장 내용", []),
        (None, []),              # ← 완전히 빈 장
        ("셋째 장 내용", []),
        (None, []),              # ← 완전히 빈 장 (마지막)
    ])

    result = PptExtractor().extract_structured_data(path)

    numbers = [s["slide_number"] for s in result["slides"]]
    assert numbers == [1, 2, 3, 4], f"장이 사라지거나 번호가 밀렸다: {numbers}"
    by_number = {s["slide_number"]: s["content"] for s in result["slides"]}
    assert "첫 장 내용" in by_number[1]
    assert by_number[2] == ""
    assert "셋째 장 내용" in by_number[3]
    assert by_number[4] == ""


def test_vision_failure_keeps_the_page(tmp_path, monkeypatch):
    """이미지 인식이 전부 실패해도 그 장은 빈 내용으로 남아야 한다(장 자체가 사라지면 안 됨)."""
    extractor = PptExtractor()
    monkeypatch.setattr(
        extractor.image_text_extractor, "extract_text_from_image",
        lambda *a, **k: (_ for _ in ()).throw(RuntimeError("vision down")))
    # 실제 API로 새면 안 된다 — 장면 묘사(3차)도 죽은 상황을 가정한다.
    monkeypatch.setattr(extractor.image_text_extractor, "describe_scene", lambda *a, **k: "")
    extractor.claude_ocr.use_fallback = True
    path = _deck(tmp_path, [
        ("첫 장 내용", []),
        (None, [(800, 600)]),    # ← 이미지뿐인데 인식이 죽는 장
    ])

    result = extractor.extract_structured_data(path)

    numbers = [s["slide_number"] for s in result["slides"]]
    assert numbers == [1, 2], f"인식 실패한 장이 사라졌다: {numbers}"
    assert result["slides"][1]["content"] == ""


def test_pdf_empty_pages_are_kept(tmp_path):
    """PDF도 같은 규칙 — 빈 페이지를 버리면 장수가 어긋난다."""
    import pypdf
    from utils import pdf_extractor

    writer = pypdf.PdfWriter()
    for _ in range(3):
        writer.add_blank_page(width=720, height=540)   # 텍스트 없는 페이지 3장
    path = tmp_path / "empty.pdf"
    with open(path, "wb") as f:
        writer.write(f)

    result = pdf_extractor.extract_structured_data(str(path))

    numbers = [s["slide_number"] for s in result["slides"]]
    assert numbers == [1, 2, 3], f"빈 페이지가 사라졌다: {numbers}"
    assert all(s["content"] == "" for s in result["slides"])


def test_template_images_do_not_hog_ocr_slots(tmp_path):
    """모든 장에 반복되는 배경/배너가 넓이 상위 자리를 차지하면 내용 이미지가 안 읽힌다.

    실측(2026-08-18, 이미지형 14장): 전 장에 깔린 1761×773 배경이 상위 3자리를 먹어서
    4개 장이 원문 0자가 됐다 — 반복 이미지는 템플릿으로 보고 후보에서 뺀다.
    """
    # 같은 (800,600) 배경이 4장 전부에 깔려 있고, 각 장에 서로 다른 내용 이미지가 있다.
    slides = [(None, [(800, 600), (300, 250)]),
              (None, [(800, 600), (310, 250)]),
              (None, [(800, 600), (320, 250)]),
              (None, [(800, 600), (330, 250)])]
    recorder = _Recorder()
    result = _extract(tmp_path, slides, recorder)

    read_sizes = set(recorder.calls)
    assert (800, 600) not in read_sizes, "템플릿 배경을 읽었다 (내용 이미지 자리를 뺏는다)"
    assert {(300, 250), (310, 250), (320, 250), (330, 250)} <= read_sizes, \
        "내용 이미지가 읽히지 않았다"
    # 내용도 정상으로 붙었는지
    contents = {s["slide_number"]: s["content"] for s in result["slides"]}
    assert contents[1] == "이미지300x250"


def test_template_only_slide_still_reads_something(tmp_path):
    """표지처럼 배경(=템플릿)뿐인 장은, 아예 안 읽는 것보다 배경이라도 읽는 게 낫다."""
    slides = [(None, [(800, 600)]),
              (None, [(800, 600), (300, 250)]),
              (None, [(800, 600), (310, 250)])]
    recorder = _Recorder()
    _extract(tmp_path, slides, recorder)

    # (800,600)은 3장 반복이라 템플릿이지만, 1장은 그것뿐이므로 거기서는 읽어야 한다.
    assert (800, 600) in set(recorder.calls), "템플릿뿐인 장에서 아무것도 안 읽었다"


class _SelectiveRecorder(_Recorder):
    """특정 크기 목록만 텍스트를 돌려주고 나머지는 빈 문자열(글자 없음)을 돌려준다."""

    def __init__(self, text_sizes, **kwargs):
        super().__init__(**kwargs)
        self.text_sizes = set(text_sizes)

    def extract_text_from_image(self, blob, context_hint, content_type):
        result = super().extract_text_from_image(blob, context_hint, content_type)
        import io as _io
        from PIL import Image as _Image
        size = _Image.open(_io.BytesIO(blob)).size
        return result if size in self.text_sizes else ""


def test_second_round_reads_more_images_when_top3_are_blank(tmp_path):
    """상위 3장이 전부 무늬/배경이면 예비 후보를 마저 읽어야 한다.

    실측(2026-08-19, 이미지형 2·3·13장): 말풍선 텍스트("반갑습니다 진순입니다:)")가 든
    캐릭터 카드가 넓이 4번째라 1차에서 안 읽혀 원문 0자가 됐다.
    """
    # 큰 3장은 글자 없음, 4번째(300x250)에만 글자가 있다.
    slides = [(None, [(900, 700), (800, 650), (700, 600), (300, 250)])]
    recorder = _SelectiveRecorder(text_sizes=[(300, 250)])
    result = _extract(tmp_path, slides, recorder)

    assert (300, 250) in set(recorder.calls), "2차 라운드가 예비 후보를 읽지 않았다"
    assert result["slides"][0]["content"] == "이미지300x250"


def test_no_second_round_when_first_round_found_text(tmp_path):
    """1차에서 글자를 찾은 장은 예비 후보를 읽지 않는다 — 유료 호출이다."""
    slides = [(None, [(900, 700), (800, 650), (700, 600), (300, 250)])]
    recorder = _SelectiveRecorder(text_sizes=[(900, 700)])   # 첫 장부터 글자 있음
    _extract(tmp_path, slides, recorder)

    assert (300, 250) not in set(recorder.calls), "글자를 이미 찾았는데 예비까지 읽었다"
    assert len(recorder.calls) == 3


# ---------------------------------------------------------------------------
# 3차(장면 묘사) + Claude OCR 폴백 (2026-08-19)
# HCX-005 비전이 큰 일반 폰트도 못 읽는다는 실험 결과에 따른 두 안전망.
# ---------------------------------------------------------------------------


class _DescribingRecorder(_Recorder):
    """글자는 하나도 못 읽지만 장면 설명은 돌려주는 비전 흉내."""

    def extract_text_from_image(self, blob, context_hint, content_type):
        super().extract_text_from_image(blob, context_hint, content_type)
        return ""

    def describe_scene(self, blob, content_type):
        return "노란 옷을 입은 캐릭터가 인사하는 장면"


def test_scene_description_is_labeled_for_textless_slide(tmp_path):
    """글자를 못 읽은 장은 [화면 묘사] 라벨로 장면 설명이 실려야 한다.

    라벨 없이 넣으면 생성기가 묘사를 슬라이드 원문으로 믿는다(과거 환각 사고).
    """
    extractor = PptExtractor()
    extractor.image_text_extractor = _DescribingRecorder()
    extractor.claude_ocr.use_fallback = True          # 키 없음 가정
    result = extractor.extract_structured_data(_deck(tmp_path, [(None, [(600, 400)])]))

    content = result["slides"][0]["content"]
    assert content.startswith("[화면 묘사]"), content
    assert "인사하는 장면" in content


def test_scene_description_not_called_when_text_was_read(tmp_path):
    """글자를 읽은 장은 장면 설명을 부르지 않는다 — 유료 호출이고 오염 위험이다."""
    calls = []

    class _TextRecorder(_Recorder):
        def describe_scene(self, blob, content_type):
            calls.append(1)
            return "장면"

    extractor = PptExtractor()
    extractor.image_text_extractor = _TextRecorder()
    extractor.claude_ocr.use_fallback = True
    result = extractor.extract_structured_data(_deck(tmp_path, [(None, [(600, 400)])]))

    assert calls == []
    assert result["slides"][0]["content"] == "이미지600x400"


def test_claude_ocr_rescues_slides_hcx_cannot_read(tmp_path):
    """HCX가 0자를 준 장은 Claude OCR이 읽고, 성공하면 장면 묘사 없이 진짜 글자가 실린다."""
    scene_calls = []

    class _BlindRecorder(_Recorder):
        def extract_text_from_image(self, blob, context_hint, content_type):
            super().extract_text_from_image(blob, context_hint, content_type)
            return ""
        def describe_scene(self, blob, content_type):
            scene_calls.append(1)
            return "장면"

    class _FakeClaude:
        use_fallback = False
        def extract_text_from_image(self, blob, content_type):
            return "반갑습니다! 진순 입니다:)"

    extractor = PptExtractor()
    extractor.image_text_extractor = _BlindRecorder()
    extractor.claude_ocr = _FakeClaude()
    result = extractor.extract_structured_data(_deck(tmp_path, [(None, [(600, 400)])]))

    assert "반갑습니다! 진순 입니다:)" in result["slides"][0]["content"]
    assert scene_calls == [], "Claude가 읽었는데 장면 묘사까지 불렀다"


def test_claude_ocr_skipped_without_key(tmp_path):
    """키가 없으면 Claude 경로는 조용히 건너뛴다 — 기동/업로드가 죽으면 안 된다."""
    calls = []

    class _FakeClaude:
        use_fallback = True
        def extract_text_from_image(self, blob, content_type):
            calls.append(1)
            return "안 불려야 함"

    extractor = PptExtractor()
    extractor.image_text_extractor = _DescribingRecorder()
    extractor.claude_ocr = _FakeClaude()
    result = extractor.extract_structured_data(_deck(tmp_path, [(None, [(600, 400)])]))

    assert calls == []
    assert result["slides"][0]["content"].startswith("[화면 묘사]")
