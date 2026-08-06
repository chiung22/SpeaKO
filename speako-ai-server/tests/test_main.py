import io
import os

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches
import docx

import main
from main import app
from db import models
from utils import pdf_extractor
from clova.full_generation import generator as full_gen_module
from clova.partial_generation import generator as partial_gen_module

client = TestClient(app)


class _FakeResponse:
    def __init__(self, payload, status_code=200, text=""):
        self._payload = payload
        # 클라이언트가 status_code를 직접 보고 4xx/5xx면 응답 본문을 로그로 남긴다.
        self.status_code = status_code
        self.text = text

    def raise_for_status(self):
        pass

    def json(self):
        return self._payload


def _make_pptx_bytes(slide_texts):
    """python-pptx로 텍스트박스가 있는 최소 pptx를 즉석에서 만들어 바이트로 반환한다."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for text in slide_texts:
        slide = prs.slides.add_slide(blank_layout)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
        box.text_frame.text = text

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def _make_docx_bytes(paragraphs):
    """python-docx로 최소 docx를 즉석에서 만들어 바이트로 반환한다."""
    document = docx.Document()
    for text in paragraphs:
        document.add_paragraph(text)

    buffer = io.BytesIO()
    document.save(buffer)
    buffer.seek(0)
    return buffer


def _create_project(db_session_factory, slides, script_map=None):
    """(slide_number -> source_content) 프로젝트를 DB에 직접 심는다. script_map이 있으면 슬라이드에 대본도 채운다."""
    script_map = script_map or {}
    db = db_session_factory()
    try:
        project = models.Project(name="테스트 프로젝트", filename="test.pptx", topic="테스트 주제", keywords=[])
        project.slides = [
            models.Slide(slide_number=num, source_content=content, script=script_map.get(num))
            for num, content in slides
        ]
        db.add(project)
        db.commit()
        db.refresh(project)
        return project.id
    finally:
        db.close()


def test_root_returns_ok():
    response = client.get("/")
    assert response.status_code == 200
    assert response.json()["message"]


def test_api_key_disabled_by_default_in_dev(db_session_factory):
    # .env에 SPEAKO_API_KEY가 플레이스홀더 상태(로컬 개발 기본값)면 인증 없이 /api/*가 열려 있어야 한다.
    response = client.get("/api/projects")
    assert response.status_code == 200


def test_api_rejects_missing_or_wrong_key_when_enabled(monkeypatch, db_session_factory):
    monkeypatch.setattr(main, "_AUTH_ENABLED", True)
    monkeypatch.setattr(main, "SPEAKO_API_KEY", "correct-secret")

    no_header = client.get("/api/projects")
    assert no_header.status_code == 401

    wrong_key = client.get("/api/projects", headers={"X-API-Key": "wrong"})
    assert wrong_key.status_code == 401

    right_key = client.get("/api/projects", headers={"X-API-Key": "correct-secret"})
    assert right_key.status_code == 200


def test_root_bypasses_api_key_even_when_auth_enabled(monkeypatch):
    monkeypatch.setattr(main, "_AUTH_ENABLED", True)
    monkeypatch.setattr(main, "SPEAKO_API_KEY", "correct-secret")

    response = client.get("/")
    assert response.status_code == 200


class _FakePdfPage:
    def __init__(self, text):
        self._text = text

    def extract_text(self):
        return self._text


class _FakePdfReader:
    def __init__(self, *_args, **_kwargs):
        # 두 번째 페이지는 텍스트가 없는 캡처 슬라이드를 흉내낸다 — 결과에서 제외되어야 한다.
        self.pages = [_FakePdfPage("첫 페이지 내용입니다"), _FakePdfPage("")]


def test_create_project_from_pptx_persists_slides(db_session_factory):
    pptx_bytes = _make_pptx_bytes(["발표 주제입니다", "두 번째 슬라이드 내용"])

    response = client.post(
        "/api/projects",
        files={"file": ("slides.pptx", pptx_bytes, "application/octet-stream")},
        data={"project_name": "내 발표"},
    )
    assert response.status_code == 200
    body = response.json()
    project_id = body["project_id"]
    assert project_id

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        assert project.name == "내 발표"
        assert len(project.slides) == 2
        assert project.slides[0].source_content == "발표 주제입니다"
    finally:
        db.close()


def test_create_project_from_pdf_persists_slides(monkeypatch, db_session_factory):
    monkeypatch.setattr(pdf_extractor.pypdf, "PdfReader", lambda *_a, **_k: _FakePdfReader())

    fake_file = io.BytesIO(b"%PDF-1.4 fake pdf bytes")
    response = client.post(
        "/api/projects",
        files={"file": ("slides.pdf", fake_file, "application/pdf")},
    )
    assert response.status_code == 200
    project_id = response.json()["project_id"]

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        assert len(project.slides) == 1  # 빈 텍스트 페이지는 제외됨
        assert project.slides[0].source_content == "첫 페이지 내용입니다"
    finally:
        db.close()


def test_create_project_from_topic_and_outline_without_file(db_session_factory):
    response = client.post(
        "/api/projects",
        data={"topic": "AI 기반 발표 코칭 서비스 기획", "outline": "1. 발표 내용\n2. 설명 대상"},
    )
    assert response.status_code == 200
    project_id = response.json()["project_id"]

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        assert project.topic == "AI 기반 발표 코칭 서비스 기획"
        assert len(project.slides) == 1
        assert "발표 내용" in project.slides[0].source_content
        assert project.slides[0].script is None  # 아직 대본 생성 전
    finally:
        db.close()


def test_create_project_from_script_text_skips_generation(db_session_factory):
    response = client.post(
        "/api/projects",
        data={"script_text": "이미 완성된 발표 대본입니다."},
    )
    assert response.status_code == 200
    project_id = response.json()["project_id"]

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        assert project.slides[0].script == "이미 완성된 발표 대본입니다."  # 생성 없이 바로 대본으로 저장됨
    finally:
        db.close()


def test_create_project_coaching_mode_from_docx_skips_generation(db_session_factory):
    docx_bytes = _make_docx_bytes(["안녕하세요 발표를 시작하겠습니다.", "오늘 다룰 내용은 다음과 같습니다."])

    response = client.post(
        "/api/projects",
        data={"mode": "coaching", "project_name": "코칭 테스트"},
        files={"file": ("my_script.docx", docx_bytes, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
    )
    assert response.status_code == 200
    project_id = response.json()["project_id"]

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        assert "안녕하세요 발표를 시작하겠습니다." in project.slides[0].script
        assert "오늘 다룰 내용은" in project.slides[0].script
    finally:
        db.close()


def test_create_project_coaching_mode_from_txt_skips_generation(db_session_factory):
    txt_bytes = io.BytesIO("붙여넣기 대신 파일로 올린 대본입니다.".encode("utf-8"))

    response = client.post(
        "/api/projects",
        data={"mode": "coaching"},
        files={"file": ("script.txt", txt_bytes, "text/plain")},
    )
    assert response.status_code == 200
    project_id = response.json()["project_id"]

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        assert project.slides[0].script == "붙여넣기 대신 파일로 올린 대본입니다."
    finally:
        db.close()


def test_create_project_coaching_mode_from_pdf_joins_all_pages(monkeypatch, db_session_factory):
    monkeypatch.setattr(pdf_extractor.pypdf, "PdfReader", lambda *_a, **_k: _FakePdfReader())

    fake_file = io.BytesIO(b"%PDF-1.4 fake pdf bytes")
    response = client.post(
        "/api/projects",
        data={"mode": "coaching"},
        files={"file": ("script.pdf", fake_file, "application/pdf")},
    )
    assert response.status_code == 200
    project_id = response.json()["project_id"]

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        assert project.slides[0].script == "첫 페이지 내용입니다"  # _FakePdfReader의 두 번째 페이지는 빈 텍스트라 제외됨
    finally:
        db.close()


def test_create_project_coaching_mode_rejects_pptx():
    # coaching 모드는 완성된 문서(DOCX/TXT/PDF)만 받는다 — 슬라이드 덱(PPTX)은 대상이 아니다.
    pptx_bytes = _make_pptx_bytes(["내용"])
    response = client.post(
        "/api/projects",
        data={"mode": "coaching"},
        files={"file": ("slides.pptx", pptx_bytes, "application/octet-stream")},
    )
    assert response.status_code == 415


def test_create_project_requires_some_input():
    response = client.post("/api/projects", data={})
    assert response.status_code == 422


def test_create_project_rejects_invalid_file():
    fake_file = io.BytesIO(b"not a real pptx file")
    response = client.post(
        "/api/projects",
        files={"file": ("broken.pptx", fake_file, "application/octet-stream")},
    )
    assert response.status_code == 422


def test_create_project_rejects_wrong_extension():
    fake_file = io.BytesIO(b"not a pptx at all")
    response = client.post(
        "/api/projects",
        files={"file": ("notes.txt", fake_file, "text/plain")},
    )
    assert response.status_code == 415


def test_create_project_rejects_oversized_file(monkeypatch):
    # 매 요청마다 20MB 페이로드를 만들지 않도록, 제한 값만 낮춰서 검증한다.
    monkeypatch.setattr(main, "MAX_PPT_SIZE_BYTES", 10)
    fake_file = io.BytesIO(b"x" * 100)
    response = client.post(
        "/api/projects",
        files={"file": ("slides.pptx", fake_file, "application/octet-stream")},
    )
    assert response.status_code == 413


def _generate_full_and_wait(payload):
    """비동기 대본 생성: 접수번호(job_id)를 받고 완료될 때까지 상태를 폴링해 최종 job 응답을 돌려준다.
    (테스트에서는 작업이 인라인으로 즉시 실행되므로 사실상 한 번의 조회로 끝난다.)"""
    start = client.post("/api/script/full", json=payload)
    assert start.status_code == 202, start.text
    job_id = start.json()["job_id"]
    for _ in range(50):
        res = client.get(f"/api/script/jobs/{job_id}")
        assert res.status_code == 200
        body = res.json()
        if body["status"] != "processing":
            return body
    raise AssertionError("작업이 완료되지 않았습니다(폴링 초과).")


def test_script_full_requires_existing_project():
    # 프로젝트 존재 확인은 job을 만들기 전에 하므로 POST가 곧바로 404를 낸다.
    response = client.post(
        "/api/script/full",
        json={"project_id": 9999, "presentation_time": 1, "style": "격식체"},
    )
    assert response.status_code == 404


def test_script_full_returns_job_id_immediately(monkeypatch, db_session_factory):
    """생성 요청은 접수번호(job_id)를 즉시 202로 돌려준다(요청을 붙잡지 않는다)."""
    project_id = _create_project(db_session_factory, [(1, "내용")])
    monkeypatch.setattr(main.full_generator, "use_fallback", False)
    monkeypatch.setattr(full_gen_module.requests, "post",
                        lambda *a, **k: _FakeResponse({"result": {"message": {"content": "한 문장 대본입니다."}}}))

    response = client.post(
        "/api/script/full",
        json={"project_id": project_id, "presentation_time": 1, "style": "격식체"},
    )
    assert response.status_code == 202
    body = response.json()
    assert body["job_id"]
    assert body["status"] == "processing"


def test_script_job_404_for_unknown_id():
    assert client.get("/api/script/jobs/does-not-exist").status_code == 404


def test_script_full_fails_without_api_key(monkeypatch, db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "테스트 슬라이드 내용")])
    # 키가 없으면 use_fallback=True가 되어 네트워크 호출 없이 None을 반환한다.
    # 이제 생성은 백그라운드 작업이므로, 접수는 202로 받고 작업 상태가 'failed'로 끝나야 한다.
    # (실제 네트워크를 때리지 않으므로 CI/오프라인에서도 결정적으로 동작한다)
    monkeypatch.setattr(main.full_generator, "use_fallback", True)
    body = _generate_full_and_wait({"project_id": project_id, "presentation_time": 1, "style": "격식체"})
    assert body["status"] == "failed"
    assert body["error"]


def test_script_full_parses_real_world_toon_variant_and_saves_to_slides(monkeypatch, db_session_factory):
    # HCX-005는 v3 chat-completions 전용이며, 실제 응답은 프롬프트의 헤더+행 구조를
    # 정확히 지키지 않고 slides[N]{...}를 슬라이드마다 반복하기도 한다.
    # 네트워크 호출 없이, 그 변형된 실제 응답 형태를 그대로 파싱하고 각 슬라이드에 저장할 수 있는지 검증한다.
    project_id = _create_project(db_session_factory, [(1, "메타버스 개념"), (2, "시장 규모")])

    raw_toon = (
        "slides[2]{1,메타버스의 개념에 대해 설명하겠습니다.}\n\n"
        "slides[2]{2,이제 시장 규모에 대해 알아보겠습니다.}\n\n"
        "이상 발표를 마치겠습니다. 감사합니다."
    )
    fake_payload = {"result": {"message": {"content": raw_toon}}}
    # CI에는 .env가 없어 HCX 키가 미설정이면 use_fallback=True가 되므로, 모킹된 네트워크 경로를
    # 타도록 강제로 False로 둔다. (이 테스트는 파싱 로직 검증이지 키 유무 검증이 아니다)
    monkeypatch.setattr(main.full_generator, "use_fallback", False)
    monkeypatch.setattr(full_gen_module.requests, "post", lambda *a, **k: _FakeResponse(fake_payload))

    body = _generate_full_and_wait({"project_id": project_id, "presentation_time": 1, "style": "격식체"})
    assert body["status"] == "completed"
    slides = body["data"]["slides"]
    assert [s["slide_number"] for s in slides] == ["1", "2"]
    assert "메타버스의 개념" in slides[0]["script"]

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        saved_scripts = {s.slide_number: s.script for s in project.slides}
        assert "메타버스의 개념" in saved_scripts[1]
        assert "시장 규모" in saved_scripts[2]
    finally:
        db.close()


def test_script_full_creates_new_slides_when_model_splits_more_than_source(monkeypatch, db_session_factory):
    # topic/outline만으로 만든 프로젝트는 원본 슬라이드가 1개뿐이지만, 모델이 여러 슬라이드로 쪼개 생성할 수 있다.
    # 기존에 없는 슬라이드 번호라도 결과가 유실되지 않고 새로 저장되어야 한다.
    project_id = _create_project(db_session_factory, [(1, "발표 주제: SpeaKO\n목차: 1. 문제 2. 기능 3. 효과")])

    raw_toon = (
        "slides[3]{slide_number,script}:\n"
        " 1,안녕하세요 SpeaKO를 소개합니다.\n"
        " 2,먼저 문제 정의입니다.\n"
        " 3,핵심 기능을 살펴보겠습니다."
    )
    fake_payload = {"result": {"message": {"content": raw_toon}}}
    monkeypatch.setattr(main.full_generator, "use_fallback", False)
    monkeypatch.setattr(full_gen_module.requests, "post", lambda *a, **k: _FakeResponse(fake_payload))

    body = _generate_full_and_wait({"project_id": project_id, "presentation_time": 3, "style": "편안한 말투"})
    assert body["status"] == "completed"

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        assert len(project.slides) == 3
        saved_scripts = {s.slide_number: s.script for s in project.slides}
        assert "문제 정의" in saved_scripts[2]
        assert "핵심 기능" in saved_scripts[3]
    finally:
        db.close()


def test_script_partial_rejects_invalid_style():
    # style은 "격식체"/"편안한 말투" 둘 중 하나만 허용해야 한다 (프로젝트 존재 여부와 무관하게 422).
    response = client.post(
        "/api/script/partial",
        json={"project_id": 1, "target_slide": 3, "style": "반말"},
    )
    assert response.status_code == 422


def test_script_partial_requires_existing_slide(db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "내용")], script_map={1: "기존 대본"})
    response = client.post(
        "/api/script/partial",
        json={"project_id": project_id, "target_slide": 99, "style": "격식체"},
    )
    assert response.status_code == 404


def test_script_partial_requires_generated_script_first(db_session_factory):
    project_id = _create_project(db_session_factory, [(3, "내용")])  # script 아직 없음
    response = client.post(
        "/api/script/partial",
        json={"project_id": project_id, "target_slide": 3, "style": "격식체"},
    )
    assert response.status_code == 422


def test_script_partial_parses_toon_and_updates_slide(monkeypatch, db_session_factory):
    # 원본 대본을 클라이언트가 다시 안 보내도, DB에 저장된 걸 그대로 써서 재생성되고
    # TOON 응답에서 slide_number/script를 구조화된 데이터로 뽑아내야 한다.
    project_id = _create_project(db_session_factory, [(3, "내용")], script_map={3: "기존 대본입니다."})

    raw_toon = "slides[1]{slide_number,script}:\n 3,다시 쓴 세 번째 슬라이드 대본입니다."
    fake_payload = {"result": {"message": {"content": raw_toon}}}
    monkeypatch.setattr(main.partial_generator, "use_fallback", False)
    monkeypatch.setattr(partial_gen_module.requests, "post", lambda *a, **k: _FakeResponse(fake_payload))

    response = client.post(
        "/api/script/partial",
        json={"project_id": project_id, "target_slide": 3, "style": "편안한 말투"},
    )
    assert response.status_code == 200
    data = response.json()["data"]
    assert data["slide_number"] == "3"
    assert "다시 쓴 세 번째 슬라이드" in data["script"]

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        assert "다시 쓴 세 번째 슬라이드" in project.slides[0].script
    finally:
        db.close()


def test_script_partial_fails_without_api_key(monkeypatch, db_session_factory):
    project_id = _create_project(db_session_factory, [(3, "내용")], script_map={3: "기존 대본입니다."})
    # 키가 없으면 use_fallback=True → 네트워크 없이 None 반환 → 502 (오프라인에서도 결정적)
    monkeypatch.setattr(main.partial_generator, "use_fallback", True)
    response = client.post(
        "/api/script/partial",
        json={
            "project_id": project_id,
            "target_slide": 3,
            "style": "격식체",
            "extra_requirement": "속도감 있게 해줘",
        },
    )
    assert response.status_code == 502


def test_analysis_words_requires_generated_script(db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "내용")])  # script 없음
    response = client.post("/api/analysis/words", json={"project_id": project_id})
    assert response.status_code == 422


def test_analysis_words_uses_kiwi_when_etri_unavailable(db_session_factory):
    # ETRI 키가 없어도, Kiwi 로컬 형태소 분석으로 실제 대본에서 명사/외국어를 뽑아 G2P 변환까지 성공해야 한다.
    project_id = _create_project(
        db_session_factory, [(1, "내용")], script_map={1: "메타버스와 인프라 구축의 특징을 살펴봅시다."}
    )
    response = client.post("/api/analysis/words", json={"project_id": project_id})
    assert response.status_code == 200
    body = response.json()
    assert body["success"] is True
    words = [w["word"] for w in body["data"]["words"]]
    assert len(words) > 0
    assert set(body["data"]["summary"].keys()) == {"장단음", "연음", "표기-발음불일치"}

    # "조사를 떼고 명사만 뽑았는가"는 추출 단계의 책임이므로 그 계층에서 확인한다.
    # API 응답은 발음상 주의할 게 없는 단어를 걸러내므로("메타버스"는 철자=발음이고 장단음도
    # 아니라 빠진다) 여기서 특정 단어를 기대하면 걸러내기 규칙과 함께 깨진다.
    extracted = main.kiwi_analyzer.extract_difficult_words("메타버스와 인프라 구축의 특징을 살펴봅시다.")
    assert "메타버스" in extracted, "Kiwi가 조사를 떼고 명사를 온전히 뽑지 못했습니다"


def test_analysis_words_classifies_into_categories_and_persists(monkeypatch, db_session_factory):
    # 장단음(stdict 모킹)/연음(구조상 판정)/표기-발음불일치(나머지) 3분류가 실제로 나뉘고, DB에도 카테고리가 저장돼야 한다.
    project_id = _create_project(
        db_session_factory, [(1, "내용")], script_map={1: "대본 내용은 분석에 안 쓰인다."}
    )
    monkeypatch.setattr(
        main.etri_analyzer, "extract_difficult_words", lambda script_text: ["특징", "밭이", "국민"]
    )
    # "특징"만 장음으로 모킹 — 우선순위상 연음/표기불일치보다 장단음 판정이 먼저 적용되는지도 같이 검증됨.
    monkeypatch.setattr(main.stdict_client, "has_long_vowel", lambda word: word == "특징")

    response = client.post("/api/analysis/words", json={"project_id": project_id})
    assert response.status_code == 200
    data = response.json()["data"]

    by_word = {w["word"]: w for w in data["words"]}
    assert by_word["특징"]["category"] == "장단음"
    assert by_word["밭이"]["category"] == "연음"
    assert by_word["국민"]["category"] == "표기-발음불일치"
    assert data["summary"] == {"장단음": 1, "연음": 1, "표기-발음불일치": 1}

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        saved_categories = {w.word: w.category for w in project.difficult_words}
        assert saved_categories == {"특징": "장단음", "밭이": "연음", "국민": "표기-발음불일치"}
    finally:
        db.close()


def test_analysis_words_long_vowel_fires_even_when_spelling_matches_pronunciation(monkeypatch, db_session_factory):
    # 회귀 방지: 장단음은 철자=발음(is_different=False)인 단어에서도 판정되어야 한다.
    # (예전엔 is_different=True일 때만 확인해서 장단음이 사실상 죽은 코드였음)
    project_id = _create_project(db_session_factory, [(1, "내용")], script_map={1: "밤이 깊었습니다."})
    # "밤"은 밤(chestnut)/밤ː(night) 동형이의어라 G2P상 철자=발음(is_different=False)이지만 장음일 수 있다.
    monkeypatch.setattr(main.etri_analyzer, "extract_difficult_words", lambda script_text: ["밤"])
    monkeypatch.setattr(main.g2p_converter, "convert_words",
                        lambda words: [{"word": "밤", "phoneme": "[밤]", "is_different": False}])
    monkeypatch.setattr(main.stdict_client, "has_long_vowel", lambda word: word == "밤")

    response = client.post("/api/analysis/words", json={"project_id": project_id})
    assert response.status_code == 200
    data = response.json()["data"]
    assert data["words"][0]["category"] == "장단음"
    assert data["summary"]["장단음"] == 1


def test_analysis_words_drops_word_with_no_pronunciation_issue(monkeypatch, db_session_factory):
    # 철자=발음이고 장단음도 아니면 분류가 없다(None) = 발음상 주의할 게 없다는 뜻이다.
    # 그런 단어는 '발음 주의 단어' 목록에서 아예 빼야 한다 — 넣으면 피그마 화면에서
    # 뱃지도 설명도 빈 줄이 된다(실측 2026-08-06: 실제 대본에서 40개 중 25개가 여기 해당).
    project_id = _create_project(db_session_factory, [(1, "내용")], script_map={1: "가구 배치."})
    monkeypatch.setattr(main.etri_analyzer, "extract_difficult_words", lambda script_text: ["가구"])
    monkeypatch.setattr(main.g2p_converter, "convert_words",
                        lambda words: [{"word": "가구", "phoneme": "[가구]", "is_different": False}])
    monkeypatch.setattr(main.stdict_client, "has_long_vowel", lambda word: False)
    monkeypatch.setattr(main.stdict_client, "long_vowel_positions", lambda word: ())

    response = client.post("/api/analysis/words", json={"project_id": project_id})
    assert response.status_code == 200
    data = response.json()["data"]
    assert data["words"] == [], "발음상 주의할 게 없는 단어가 목록에 남아 있습니다"
    assert data["summary"] == {"장단음": 0, "연음": 0, "표기-발음불일치": 0}


def test_analysis_words_requires_existing_project():
    response = client.post("/api/analysis/words", json={"project_id": 999999})
    assert response.status_code == 404


def test_script_partial_requires_existing_project():
    response = client.post(
        "/api/script/partial",
        json={"project_id": 999999, "target_slide": 1, "style": "격식체"},
    )
    assert response.status_code == 404


def test_create_project_coaching_mode_rejects_corrupt_docx():
    # 손상된(내용이 docx가 아닌) 파일은 raw 500이 아니라 깨끗한 422로 알려야 한다.
    fake_file = io.BytesIO(b"this is not a real docx zip container")
    response = client.post(
        "/api/projects",
        data={"mode": "coaching"},
        files={"file": ("broken.docx", fake_file, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
    )
    assert response.status_code == 422


def test_evaluation_audio_rejects_wrong_extension():
    # 허용 목록(WAV/MP3/M4A/WEBM) 밖 확장자(예: OGG)는 여전히 거부해야 한다.
    fake_file = io.BytesIO(b"not an audio file")
    response = client.post(
        "/api/evaluation/audio",
        data={"project_id": "1", "reference_text": "테스트 문장입니다."},
        files={"audio_file": ("clip.ogg", fake_file, "audio/ogg")},
    )
    assert response.status_code == 415


def test_evaluation_audio_accepts_browser_webm(monkeypatch, db_session_factory):
    # 브라우저 MediaRecorder 기본 포맷(webm)을 받아 ffmpeg 변환 후 평가해야 한다.
    # (프론트가 녹음한 걸 그대로 던질 수 있어야 함 — webm이 415로 막히면 안 됨)
    project_id = _create_project(db_session_factory, [(1, "내용")], script_map={1: "테스트 문장입니다."})

    converted = []

    def _fake_convert(input_path, output_path):
        converted.append((input_path, output_path))
        with open(output_path, "wb") as f:
            f.write(b"fake wav bytes")
        return True

    monkeypatch.setattr(main.audio_converter, "convert_to_wav", _fake_convert)
    monkeypatch.setattr(
        main.azure_evaluator,
        "evaluate_audio",
        lambda audio_file_path, reference_text: {
            "status": "success",
            "scores": {"accuracy": 88.0, "fluency": 88.0, "completeness": 88.0, "pronunciation_score": 88.0},
            "words_detail": [],
        },
    )

    fake_file = io.BytesIO(b"fake webm bytes")
    response = client.post(
        "/api/evaluation/audio",
        data={"project_id": str(project_id)},
        files={"audio_file": ("recording.webm", fake_file, "audio/webm")},
    )
    assert response.status_code == 200
    assert len(converted) == 1  # webm은 wav가 아니므로 반드시 변환을 거쳐야 한다


def test_evaluation_audio_converts_mp3_before_evaluating(monkeypatch, db_session_factory):
    # WAV가 아닌 파일(MP3/M4A)은 ffmpeg 변환을 거친 뒤 그 결과로 평가해야 한다.
    project_id = _create_project(db_session_factory, [(1, "내용")], script_map={1: "테스트 문장입니다."})

    converted_paths = []

    def _fake_convert(input_path, output_path):
        converted_paths.append((input_path, output_path))
        with open(output_path, "wb") as f:
            f.write(b"fake wav bytes")
        return True

    evaluated_paths = []

    def _fake_evaluate(audio_file_path, reference_text):
        evaluated_paths.append(audio_file_path)
        return {
            "status": "success",
            "overall_scores": {"accuracy": 90.0, "fluency": 90.0, "completeness": 90.0, "pronunciation_score": 90.0},
            "words_detail": [],
        }

    monkeypatch.setattr(main.audio_converter, "convert_to_wav", _fake_convert)
    monkeypatch.setattr(main.azure_evaluator, "evaluate_audio", _fake_evaluate)

    fake_file = io.BytesIO(b"fake mp3 bytes")
    response = client.post(
        "/api/evaluation/audio",
        data={"project_id": str(project_id)},
        files={"audio_file": ("clip.mp3", fake_file, "audio/mpeg")},
    )
    assert response.status_code == 200
    assert len(converted_paths) == 1
    # 실제로 변환된 wav 경로로 평가가 이뤄져야 하고(원본 mp3 경로가 아니라), 변환 산출물은 정리되어야 한다.
    assert evaluated_paths == [converted_paths[0][1]]
    assert not os.path.exists(converted_paths[0][1])


def test_evaluation_audio_returns_502_when_conversion_fails(monkeypatch, db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "내용")], script_map={1: "테스트 문장입니다."})
    monkeypatch.setattr(main.audio_converter, "convert_to_wav", lambda *_a, **_k: False)

    fake_file = io.BytesIO(b"broken m4a bytes")
    response = client.post(
        "/api/evaluation/audio",
        data={"project_id": str(project_id)},
        files={"audio_file": ("clip.m4a", fake_file, "audio/mp4")},
    )
    assert response.status_code == 502


def test_evaluation_audio_requires_existing_project():
    fake_file = io.BytesIO(b"RIFF....WAVEfmt ")
    response = client.post(
        "/api/evaluation/audio",
        data={"project_id": "9999", "reference_text": "테스트 문장입니다."},
        files={"audio_file": ("clip.wav", fake_file, "audio/wav")},
    )
    assert response.status_code == 404


def test_evaluation_audio_returns_502_on_failure_and_does_not_save(monkeypatch, db_session_factory):
    # 다른 엔드포인트와 동일하게, 평가 실패는 200이 아닌 502로 알려야 한다.
    project_id = _create_project(db_session_factory, [(1, "내용")], script_map={1: "테스트 문장입니다."})
    monkeypatch.setattr(
        main.azure_evaluator,
        "evaluate_audio",
        lambda audio_file_path, reference_text: {"status": "error", "message": "평가 중 오류 발생: 테스트"},
    )

    fake_file = io.BytesIO(b"RIFF....WAVEfmt ")
    response = client.post(
        "/api/evaluation/audio",
        data={"project_id": str(project_id)},
        files={"audio_file": ("clip.wav", fake_file, "audio/wav")},
    )
    assert response.status_code == 502

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        assert len(project.evaluations) == 0
    finally:
        db.close()


def test_evaluation_audio_saves_history_on_success(monkeypatch, db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "내용")], script_map={1: "테스트 문장입니다."})
    # Azure는 소수 점수를 overall_scores 키로 준다. 백엔드는 0~5점으로 뭉개지 말고
    # 소수 1자리(0~100)까지 자세히 내려줘야 한다(미세한 발음 차이가 드러나게).
    fake_result = {
        "status": "success",
        "overall_scores": {"accuracy": 90.44, "fluency": 85.66, "completeness": 80.0, "pronunciation_score": 88.75},
        "words_detail": [{"word": "테스트", "accuracy_score": 72.73, "error_type": "None"}],
    }
    monkeypatch.setattr(main.azure_evaluator, "evaluate_audio", lambda audio_file_path, reference_text: fake_result)

    fake_file = io.BytesIO(b"RIFF....WAVEfmt ")
    response = client.post(
        "/api/evaluation/audio",
        data={"project_id": str(project_id)},  # reference_text 생략 → DB의 대본으로 평가
        files={"audio_file": ("clip.wav", fake_file, "audio/wav")},
    )
    assert response.status_code == 200
    body = response.json()
    assert body["evaluation_id"]

    # 응답 점수는 소수 1자리로 자세히 내려줘야 한다(0~5점으로 압축 금지). 프론트는 그대로 표시만 한다.
    scores = body["overall_scores"]
    assert scores == {"accuracy": 90.4, "fluency": 85.7, "completeness": 80.0, "pronunciation_score": 88.8}
    assert body["words_detail"][0]["accuracy_score"] == 72.7

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        assert len(project.evaluations) == 1
        # DB에도 소수점 그대로 저장(키 불일치 버그 회귀 방지 — 예전엔 overall_scores를 못 읽어 None 저장됐음)
        assert project.evaluations[0].accuracy_score == 90.4
    finally:
        db.close()


def test_list_and_get_project(db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "내용")], script_map={1: "생성된 대본"})

    list_response = client.get("/api/projects")
    assert list_response.status_code == 200
    assert any(p["id"] == project_id for p in list_response.json()["data"])

    detail_response = client.get(f"/api/projects/{project_id}")
    assert detail_response.status_code == 200
    data = detail_response.json()["data"]
    assert data["slides"][0]["script"] == "생성된 대본"


def test_get_project_returns_404_for_missing_project():
    response = client.get("/api/projects/999999")
    assert response.status_code == 404


# ── 대본 편집 저장 (PUT) — 피그마 05 결과 화면의 수동/자동 저장 ──────────────────

def test_update_slide_script_saves_edited_text(db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "내용1"), (2, "내용2")], script_map={1: "초안1", 2: "초안2"})

    response = client.put(f"/api/projects/{project_id}/slides/2", json={"script": "사용자가 직접 고친 대본"})
    assert response.status_code == 200

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        slide2 = next(s for s in project.slides if s.slide_number == 2)
        assert slide2.script == "사용자가 직접 고친 대본"
        # 다른 슬라이드는 건드리지 않는다.
        slide1 = next(s for s in project.slides if s.slide_number == 1)
        assert slide1.script == "초안1"
    finally:
        db.close()


def test_update_slide_script_allows_empty_string(db_session_factory):
    """빈 문자열로 비우는 것도 편집이다(422가 아니라 저장돼야 한다)."""
    project_id = _create_project(db_session_factory, [(1, "내용")], script_map={1: "초안"})
    response = client.put(f"/api/projects/{project_id}/slides/1", json={"script": ""})
    assert response.status_code == 200


def test_update_slide_script_404_for_missing_slide(db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "내용")])
    assert client.put(f"/api/projects/{project_id}/slides/99", json={"script": "x"}).status_code == 404
    assert client.put("/api/projects/999999/slides/1", json={"script": "x"}).status_code == 404


# ── 슬라이드 추가/삭제 (POST/DELETE) — 피그마 05-1 ─────────────────────────────

def test_add_slide_appends_at_end_when_no_position(db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "A"), (2, "B")])
    response = client.post(f"/api/projects/{project_id}/slides", json={"script": "새 대본"})
    assert response.status_code == 200
    slides = response.json()["data"]["slides"]
    assert [s["slide_number"] for s in slides] == [1, 2, 3]
    assert slides[2]["script"] == "새 대본"


def test_add_slide_inserts_at_position_and_shifts_following(db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "A"), (2, "B"), (3, "C")], script_map={1: "a", 2: "b", 3: "c"})
    response = client.post(f"/api/projects/{project_id}/slides", json={"position": 2, "script": "끼운 대본"})
    assert response.status_code == 200
    slides = response.json()["data"]["slides"]
    # 1..N 연속 유지 + 2번 자리에 새 슬라이드, 기존 b/c는 뒤로 밀린다.
    assert [s["slide_number"] for s in slides] == [1, 2, 3, 4]
    assert [s["script"] for s in slides] == ["a", "끼운 대본", "b", "c"]


def test_delete_slide_removes_and_resequences(db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "A"), (2, "B"), (3, "C")], script_map={1: "a", 2: "b", 3: "c"})
    response = client.delete(f"/api/projects/{project_id}/slides/2")
    assert response.status_code == 200
    slides = response.json()["data"]["slides"]
    # 2번을 지우면 3번이 2번으로 당겨져 1..N이 유지된다.
    assert [s["slide_number"] for s in slides] == [1, 2]
    assert [s["script"] for s in slides] == ["a", "c"]


def test_delete_last_remaining_slide_is_rejected(db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "only")])
    response = client.delete(f"/api/projects/{project_id}/slides/1")
    assert response.status_code == 422


def test_delete_slide_404_for_missing(db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "A"), (2, "B")])
    assert client.delete(f"/api/projects/{project_id}/slides/99").status_code == 404
    assert client.delete("/api/projects/999999/slides/1").status_code == 404


# ── 프로젝트(기록) 삭제 + 발표 코칭 내역 — 마이페이지 ─────────────────────────

def _add_evaluation(db_session_factory, project_id, accuracy):
    db = db_session_factory()
    try:
        db.add(models.PronunciationEvaluation(
            project_id=project_id, accuracy_score=accuracy, fluency_score=90.0,
            completeness_score=100.0, pronunciation_score=88.5, words_detail=[],
        ))
        db.commit()
    finally:
        db.close()


def test_delete_project_removes_it_and_cascades(db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "A"), (2, "B")], script_map={1: "a", 2: "b"})
    _add_evaluation(db_session_factory, project_id, 91.2)

    response = client.delete(f"/api/projects/{project_id}")
    assert response.status_code == 200
    assert response.json()["deleted_project_id"] == project_id

    # 프로젝트와 함께 슬라이드·평가가 cascade로 사라진다.
    assert client.get(f"/api/projects/{project_id}").status_code == 404
    db = db_session_factory()
    try:
        assert db.get(models.Project, project_id) is None
        assert db.query(models.Slide).filter_by(project_id=project_id).count() == 0
        assert db.query(models.PronunciationEvaluation).filter_by(project_id=project_id).count() == 0
    finally:
        db.close()


def test_delete_project_404_for_missing():
    assert client.delete("/api/projects/999999").status_code == 404


def _create_evaluation(db_session_factory, project_id, words_detail=None):
    db = db_session_factory()
    try:
        evaluation = models.PronunciationEvaluation(
            project_id=project_id, accuracy_score=87.4, fluency_score=82.1,
            completeness_score=95.0, pronunciation_score=86.0,
            words_detail=words_detail if words_detail is not None else
            [{"word": "특징을", "accuracy_score": 50.0, "error_type": "Mispronunciation"}],
        )
        db.add(evaluation)
        db.commit()
        db.refresh(evaluation)
        return evaluation.id
    finally:
        db.close()


def test_evaluation_feedback_generates_and_saves(monkeypatch, db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "내용")], script_map={1: "메타버스를 소개합니다."})
    evaluation_id = _create_evaluation(db_session_factory, project_id)

    monkeypatch.setattr(
        main.feedback_generator, "generate_feedback",
        lambda overall_scores, weak_words, script_excerpt="", strong_words=None: {
            "summary": "전반적으로 또렷합니다.", "strengths": ["속도가 일정합니다."],
            "improvements": ["받침을 끝까지 발음하세요."], "practice_tips": ["천천히 3번 읽어보세요."],
        },
    )

    response = client.post(f"/api/evaluation/{evaluation_id}/feedback")
    assert response.status_code == 200
    body = response.json()
    assert body["cached"] is False
    assert body["data"]["summary"] == "전반적으로 또렷합니다."
    # 어떤 단어를 근거로 지적했는지 함께 내려준다.
    assert body["data"]["weak_words"][0]["word"] == "특징을"

    # 조회 API에도 피드백이 실려 나온다.
    detail = client.get(f"/api/projects/{project_id}").json()["data"]
    assert detail["evaluations"][0]["feedback"]["summary"] == "전반적으로 또렷합니다."


def test_evaluation_feedback_is_cached_and_not_regenerated(monkeypatch, db_session_factory):
    """이미 만든 피드백이 있으면 HCX를 다시 부르지 않는다(불필요한 비용 방지)."""
    project_id = _create_project(db_session_factory, [(1, "내용")], script_map={1: "대본"})
    evaluation_id = _create_evaluation(db_session_factory, project_id)

    calls = []

    def fake_generate(overall_scores, weak_words, script_excerpt="", strong_words=None):
        calls.append(1)
        return {"summary": "첫 생성", "strengths": [], "improvements": [], "practice_tips": []}

    monkeypatch.setattr(main.feedback_generator, "generate_feedback", fake_generate)

    first = client.post(f"/api/evaluation/{evaluation_id}/feedback").json()
    second = client.post(f"/api/evaluation/{evaluation_id}/feedback").json()

    assert first["cached"] is False and second["cached"] is True
    assert second["data"]["summary"] == "첫 생성"
    assert len(calls) == 1, "두 번째 호출에서 HCX를 다시 부르면 안 된다"


def test_evaluation_feedback_502_when_generation_fails(monkeypatch, db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "내용")])
    evaluation_id = _create_evaluation(db_session_factory, project_id)
    monkeypatch.setattr(main.feedback_generator, "generate_feedback",
                        lambda overall_scores, weak_words, script_excerpt="", strong_words=None: None)

    assert client.post(f"/api/evaluation/{evaluation_id}/feedback").status_code == 502


def test_evaluation_feedback_404_for_missing_evaluation():
    assert client.post("/api/evaluation/999999/feedback").status_code == 404


def test_list_evaluations_returns_all_newest_first(db_session_factory):
    p1 = _create_project(db_session_factory, [(1, "A")])
    p2 = _create_project(db_session_factory, [(1, "B")])
    _add_evaluation(db_session_factory, p1, 70.5)
    _add_evaluation(db_session_factory, p2, 85.3)

    response = client.get("/api/evaluations")
    assert response.status_code == 200
    data = response.json()["data"]
    # 두 프로젝트의 평가가 프로젝트 구분 없이 한 목록으로 나온다.
    project_ids = {e["project_id"] for e in data}
    assert {p1, p2} <= project_ids
    # 각 항목은 어느 프로젝트인지(project_name)와 점수를 포함한다.
    sample = next(e for e in data if e["project_id"] == p2)
    assert sample["project_name"] == "테스트 프로젝트"
    assert sample["accuracy_score"] == 85.3


# ── CORS (배포된 프론트엔드가 브라우저에서 직접 호출할 수 있어야 함) ──────────────

def test_cors_allows_deployed_frontend_origin():
    """localhost만 허용하면 배포된 프론트(vercel)에서 호출 시 브라우저가 전부 차단한다."""
    response = client.options(
        "/api/projects",
        headers={
            "Origin": "https://speakofront.vercel.app",
            "Access-Control-Request-Method": "GET",
        },
    )
    assert response.status_code in (200, 204)
    assert response.headers.get("access-control-allow-origin") == "https://speakofront.vercel.app"


def test_cors_allows_vercel_preview_domains():
    """Vercel 프리뷰는 커밋마다 도메인이 바뀌므로 정규식으로도 허용돼야 한다."""
    origin = "https://speakofront-abc123-team.vercel.app"
    response = client.options(
        "/api/projects",
        headers={"Origin": origin, "Access-Control-Request-Method": "GET"},
    )
    assert response.headers.get("access-control-allow-origin") == origin


def test_cors_still_allows_localhost_for_dev():
    response = client.options(
        "/api/projects",
        headers={"Origin": "http://localhost:3000", "Access-Control-Request-Method": "GET"},
    )
    assert response.headers.get("access-control-allow-origin") == "http://localhost:3000"


def test_cors_origin_list_is_configurable():
    """환경변수로 배포 도메인을 바꿀 수 있어야 한다(도메인 변경 시 코드 수정 없이)."""
    assert main._parse_origins("https://a.com, https://b.com/") == ["https://a.com", "https://b.com"]
    # 비어 있으면 기본 목록으로 폴백한다.
    assert main._parse_origins("") == main.DEFAULT_ALLOWED_ORIGINS
    assert "https://speakofront.vercel.app" in main._parse_origins("")


# ── 원본 텍스트 ↔ 인식 텍스트 (피그마 Feedback Page 좌우 비교) ──────────────────

def test_evaluation_saves_and_returns_recognized_text(monkeypatch, db_session_factory):
    """점수만으로는 어디를 잘못 읽었는지 알 수 없다. Azure가 실제로 들은 문장을 함께 저장·반환해야 한다."""
    project_id = _create_project(db_session_factory, [(1, "내용")], script_map={1: "메타버스를 소개합니다."})
    fake_result = {
        "status": "success",
        "overall_scores": {"accuracy": 90.0, "fluency": 90.0, "completeness": 90.0, "pronunciation_score": 90.0},
        "recognized_text": "메타버스를 소개함니다",
        "words_detail": [{"word": "메타버스를", "accuracy_score": 90.0, "error_type": "None"}],
    }
    monkeypatch.setattr(main.azure_evaluator, "evaluate_audio", lambda audio_file_path, reference_text: fake_result)

    response = client.post(
        "/api/evaluation/audio",
        data={"project_id": str(project_id)},
        files={"audio_file": ("clip.wav", io.BytesIO(b"RIFF....WAVEfmt "), "audio/wav")},
    )
    assert response.status_code == 200
    assert response.json()["recognized_text"] == "메타버스를 소개함니다"

    # 조회 API에서 원본 대본과 인식 텍스트를 나란히 받을 수 있어야 한다.
    detail = client.get(f"/api/projects/{project_id}").json()["data"]["evaluations"][0]
    assert detail["recognized_text"] == "메타버스를 소개함니다"
    assert detail["reference_text"] == "Slide 1: 메타버스를 소개합니다."

    listed = client.get("/api/evaluations").json()["data"][0]
    assert listed["recognized_text"] == "메타버스를 소개함니다"
    assert listed["reference_text"]


# ── 슬라이드별 부분 녹음 평가 ────────────────────────────────────────────────

def _fake_eval_result():
    return {
        "status": "success",
        "overall_scores": {"accuracy": 88.0, "fluency": 85.0, "completeness": 92.0, "pronunciation_score": 87.0},
        "recognized_text": "인식된 문장",
        "words_detail": [{"word": "발전", "accuracy_score": 60.0, "error_type": "Mispronunciation"}],
    }


def test_evaluation_with_slide_number_uses_only_that_slide_script(monkeypatch, db_session_factory):
    """슬라이드별로 녹음하면 그 장 대본만 기준으로 채점해야 한다.
    전체 대본을 기준으로 잡으면 한 장만 읽었을 때 완성도가 바닥으로 나온다."""
    project_id = _create_project(
        db_session_factory, [(1, "A"), (2, "B")], script_map={1: "첫 번째 장 대본", 2: "두 번째 장 대본"}
    )
    seen = {}

    def fake_eval(audio_file_path, reference_text):
        seen["reference"] = reference_text
        return _fake_eval_result()

    monkeypatch.setattr(main.azure_evaluator, "evaluate_audio", fake_eval)

    response = client.post(
        "/api/evaluation/audio",
        data={"project_id": str(project_id), "slide_number": "2"},
        files={"audio_file": ("clip.wav", io.BytesIO(b"RIFF"), "audio/wav")},
    )
    assert response.status_code == 200
    # 전체 대본이 아니라 2번 슬라이드 대본만 기준이 돼야 한다.
    assert seen["reference"] == "두 번째 장 대본"
    assert response.json()["slide_number"] == 2


def test_evaluation_records_slide_number_in_history(monkeypatch, db_session_factory):
    """코칭 내역에서 '3번 슬라이드 87점'처럼 구분하려면 슬라이드 번호가 남아야 한다."""
    project_id = _create_project(db_session_factory, [(1, "A"), (2, "B")], script_map={1: "가", 2: "나"})
    monkeypatch.setattr(main.azure_evaluator, "evaluate_audio", lambda audio_file_path, reference_text: _fake_eval_result())

    client.post(
        "/api/evaluation/audio",
        data={"project_id": str(project_id), "slide_number": "2"},
        files={"audio_file": ("clip.wav", io.BytesIO(b"RIFF"), "audio/wav")},
    )

    listed = client.get("/api/evaluations").json()["data"][0]
    assert listed["slide_number"] == 2
    detail = client.get(f"/api/projects/{project_id}").json()["data"]["evaluations"][0]
    assert detail["slide_number"] == 2


def test_evaluation_without_slide_number_stays_null(monkeypatch, db_session_factory):
    """대본 전체를 한 번에 녹음한 경우는 슬라이드 번호가 없다(기존 동작 유지)."""
    project_id = _create_project(db_session_factory, [(1, "A")], script_map={1: "전체 대본"})
    monkeypatch.setattr(main.azure_evaluator, "evaluate_audio", lambda audio_file_path, reference_text: _fake_eval_result())

    response = client.post(
        "/api/evaluation/audio",
        data={"project_id": str(project_id)},
        files={"audio_file": ("clip.wav", io.BytesIO(b"RIFF"), "audio/wav")},
    )
    assert response.status_code == 200
    assert response.json()["slide_number"] is None


def test_evaluation_rejects_unknown_or_empty_slide(monkeypatch, db_session_factory):
    project_id = _create_project(db_session_factory, [(1, "A"), (2, "B")], script_map={1: "가"})
    monkeypatch.setattr(main.azure_evaluator, "evaluate_audio", lambda audio_file_path, reference_text: _fake_eval_result())

    # 없는 슬라이드 번호
    missing = client.post(
        "/api/evaluation/audio",
        data={"project_id": str(project_id), "slide_number": "99"},
        files={"audio_file": ("clip.wav", io.BytesIO(b"RIFF"), "audio/wav")},
    )
    assert missing.status_code == 404

    # 대본이 아직 없는 슬라이드
    empty = client.post(
        "/api/evaluation/audio",
        data={"project_id": str(project_id), "slide_number": "2"},
        files={"audio_file": ("clip.wav", io.BytesIO(b"RIFF"), "audio/wav")},
    )
    assert empty.status_code == 422


def test_reference_text_wins_over_slide_number(monkeypatch, db_session_factory):
    """reference_text를 직접 주면 그게 우선이다."""
    project_id = _create_project(db_session_factory, [(1, "A"), (2, "B")], script_map={1: "가", 2: "나"})
    seen = {}

    def fake_eval(audio_file_path, reference_text):
        seen["reference"] = reference_text
        return _fake_eval_result()

    monkeypatch.setattr(main.azure_evaluator, "evaluate_audio", fake_eval)
    client.post(
        "/api/evaluation/audio",
        data={"project_id": str(project_id), "slide_number": "2", "reference_text": "직접 준 문장"},
        files={"audio_file": ("clip.wav", io.BytesIO(b"RIFF"), "audio/wav")},
    )
    assert seen["reference"] == "직접 준 문장"


def test_partial_regeneration_survives_header_only_response(monkeypatch, db_session_factory):
    """모델이 TOON 헤더만 붙이고 본문은 평문으로 줘도 대본을 살려야 한다(실측 502 회귀 방지)."""
    project_id = _create_project(db_session_factory, [(1, "A"), (2, "B")], script_map={1: "가", 2: "나"})
    raw = "slides[2]{slide_number,script}: \n자 그럼 이제 발표를 시작해볼게요! 집중해 주세요."
    monkeypatch.setattr(main.partial_generator, "use_fallback", False)
    monkeypatch.setattr(
        partial_gen_module.requests, "post",
        lambda *a, **k: _FakeResponse({"result": {"message": {"content": raw}}}),
    )

    response = client.post(
        "/api/script/partial",
        json={"project_id": project_id, "target_slide": 2, "style": "편안한 말투"},
    )
    assert response.status_code == 200, response.text
    data = response.json()["data"]
    assert data["script"] == "자 그럼 이제 발표를 시작해볼게요! 집중해 주세요."
    # 요청자가 지정한 슬라이드에 저장돼야 한다 (모델이 매긴 번호를 믿으면 안 됨).
    assert data["slide_number"] == "2"

    db = db_session_factory()
    try:
        project = db.get(models.Project, project_id)
        saved = {s.slide_number: s.script for s in project.slides}
        assert saved[2] == "자 그럼 이제 발표를 시작해볼게요! 집중해 주세요."
        assert saved[1] == "가", "다른 슬라이드는 건드리면 안 된다"
    finally:
        db.close()


def test_partial_regeneration_includes_slide_source_content(monkeypatch, db_session_factory):
    """대상 슬라이드의 원문을 넘겨야 모델이 그 장이 무슨 내용인지 알고 다시 쓸 수 있다.
    (원문이 없으면 앞뒤 대본만 보고 지어낸다 — 특히 아직 대본이 없는 슬라이드)"""
    project_id = _create_project(
        db_session_factory, [(1, "첫 장 원문"), (2, "시장 규모와 성장률 도표")], script_map={1: "가"}
    )
    sent = {}

    def fake_post(url, headers=None, json=None, timeout=None):
        sent["prompt"] = json["messages"][-1]["content"][0]["text"]
        return _FakeResponse({"result": {"message": {"content": "다시 쓴 대본입니다."}}})

    monkeypatch.setattr(main.partial_generator, "use_fallback", False)
    monkeypatch.setattr(partial_gen_module.requests, "post", fake_post)

    response = client.post(
        "/api/script/partial",
        json={"project_id": project_id, "target_slide": 2, "style": "격식체"},
    )
    assert response.status_code == 200
    # 아직 대본이 없는 2번 슬라이드라도 원문이 프롬프트에 들어가야 한다.
    assert "시장 규모와 성장률 도표" in sent["prompt"]
    assert "대상 슬라이드 원문" in sent["prompt"]


def test_script_job_reports_missing_slides(monkeypatch, db_session_factory):
    """생성 실패한 슬라이드가 있으면 폴링 응답으로 프론트가 알 수 있어야 한다."""
    project_id = _create_project(db_session_factory, [(1, "가"), (2, "나")])

    def fake_post(url, headers=None, json=None, timeout=None):
        prompt = json["messages"][-1]["content"][0]["text"]
        content = "" if "Slide 2" in prompt else "정상 대본입니다."
        return _FakeResponse({"result": {"message": {"content": content}}})

    monkeypatch.setattr(main.full_generator, "use_fallback", False)
    monkeypatch.setattr(full_gen_module.requests, "post", fake_post)

    body = _generate_full_and_wait({"project_id": project_id, "presentation_time": 2, "style": "격식체"})
    assert body["status"] == "completed"
    assert body["data"]["missing_slide_numbers"] == ["2"]
