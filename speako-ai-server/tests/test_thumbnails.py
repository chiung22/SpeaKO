"""슬라이드 썸네일(편집 화면 왼쪽 미리보기) 검증.

브라우저에서 PPTX를 그릴 방법이 없어 서버가 PNG로 만들어 준다. 변환 자체는 LibreOffice와
pdftoppm이라는 외부 프로그램이라 테스트 환경에 없을 수 있으므로, 여기서는

- 도구가 없을 때 **업로드가 실패하지 않고** 조용히 건너뛰는지
- 상태(pending/ready/failed/unavailable/skipped)가 응답에 정확히 실리는지
- 만들어진 PNG가 엔드포인트로 나오고, 없으면 404인지
- 프로젝트를 지울 때 디스크의 썸네일도 같이 지워지는지

를 본다. 실제 변환 결과(한글이 두부로 깨지지 않는지 등)는 도커 이미지에서 따로 확인한다.
"""
import io
import os

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches

import main
from main import app
from db import models
from utils import thumbnail_generator

client = TestClient(app)


@pytest.fixture(autouse=True)
def _isolate_thumbnail_dir(monkeypatch, tmp_path):
    """테스트가 실제 data/thumbnails에 파일을 남기지 않도록 격리한다."""
    monkeypatch.setattr(thumbnail_generator, "THUMBNAIL_DIR", str(tmp_path / "thumbnails"))


def _pptx_bytes(slide_texts):
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    for text in slide_texts:
        slide = presentation.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
        box.text_frame.text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer


def _write_thumbnail(project_id, slide_number, content=b"\x89PNG\r\n\x1a\nfake"):
    directory = thumbnail_generator.project_dir(project_id)
    os.makedirs(directory, exist_ok=True)
    path = os.path.join(directory, f"{slide_number}.png")
    with open(path, "wb") as handle:
        handle.write(content)
    return path


def _project(db_session_factory, slide_count=3, thumbnail_status="ready"):
    db = db_session_factory()
    try:
        project = models.Project(name="썸네일 테스트", filename="deck.pptx", topic=None,
                                 keywords=[], thumbnail_status=thumbnail_status)
        project.slides = [
            models.Slide(slide_number=n, source_content=f"{n}번 원문", script=f"{n}번 대본")
            for n in range(1, slide_count + 1)
        ]
        db.add(project)
        db.commit()
        db.refresh(project)
        return project.id
    finally:
        db.close()


# ------------------------------------------------------------ 파일 저장/조회

def test_thumbnail_path_is_none_when_missing():
    assert thumbnail_generator.thumbnail_path(999, 1) is None


def test_available_slide_numbers_is_sorted_numerically():
    """10장 이상이면 문자열 정렬로는 1, 10, 2 순이 된다. 화면 순서가 뒤집히면 안 된다."""
    for number in (1, 2, 10, 11, 3):
        _write_thumbnail(7, number)

    assert thumbnail_generator.available_slide_numbers(7) == [1, 2, 3, 10, 11]


def test_available_slide_numbers_ignores_junk_files():
    _write_thumbnail(8, 1)
    directory = thumbnail_generator.project_dir(8)
    open(os.path.join(directory, "notes.txt"), "w").close()
    open(os.path.join(directory, "cover.png"), "w").close()  # 숫자가 아닌 이름

    assert thumbnail_generator.available_slide_numbers(8) == [1]


def test_clear_project_removes_the_directory():
    _write_thumbnail(9, 1)
    assert thumbnail_generator.available_slide_numbers(9) == [1]

    thumbnail_generator.clear_project(9)

    assert thumbnail_generator.available_slide_numbers(9) == []


def test_clear_project_is_safe_when_nothing_exists():
    thumbnail_generator.clear_project(12345)  # 예외가 나면 삭제 API가 500이 된다


# ------------------------------------------------------------ 도구가 없을 때

def test_generate_reports_unavailable_without_tools(monkeypatch):
    """로컬 개발 환경엔 LibreOffice가 없다. 그래도 업로드는 성공해야 한다."""
    monkeypatch.setattr(thumbnail_generator, "is_available", lambda: False)

    result = thumbnail_generator.generate_for_project(1, "없는파일.pptx")

    assert result["status"] == "unavailable"
    assert result["count"] == 0


def test_generate_reports_failure_when_source_missing(monkeypatch):
    monkeypatch.setattr(thumbnail_generator, "is_available", lambda: True)

    result = thumbnail_generator.generate_for_project(1, "존재하지_않는_파일.pptx")

    assert result["status"] == "failed"


# ------------------------------------------------------------ 엔드포인트

def test_thumbnail_endpoint_returns_png(db_session_factory):
    project_id = _project(db_session_factory)
    _write_thumbnail(project_id, 2)

    response = client.get(f"/api/projects/{project_id}/slides/2/thumbnail")

    assert response.status_code == 200
    assert response.headers["content-type"] == "image/png"
    assert response.content.startswith(b"\x89PNG")


def test_thumbnail_endpoint_404_when_not_generated(db_session_factory):
    project_id = _project(db_session_factory, thumbnail_status="pending")

    response = client.get(f"/api/projects/{project_id}/slides/1/thumbnail")

    assert response.status_code == 404
    # 왜 없는지(아직 만드는 중인지, 실패했는지) 알려줘야 프론트가 재시도 여부를 정한다
    assert "pending" in response.json()["detail"]


def test_thumbnail_endpoint_404_for_unknown_project():
    assert client.get("/api/projects/999999/slides/1/thumbnail").status_code == 404


def test_project_detail_marks_which_slides_have_thumbnails(db_session_factory):
    """프론트는 has_thumbnail이 true인 것만 요청해야 깨진 이미지가 안 뜬다."""
    project_id = _project(db_session_factory, slide_count=3)
    _write_thumbnail(project_id, 1)
    _write_thumbnail(project_id, 3)

    body = client.get(f"/api/projects/{project_id}").json()["data"]

    assert body["thumbnail_status"] == "ready"
    assert [s["has_thumbnail"] for s in body["slides"]] == [True, False, True]


def test_project_list_includes_thumbnail_status(db_session_factory):
    _project(db_session_factory, thumbnail_status="pending")

    listed = client.get("/api/projects").json()["data"][0]

    assert listed["thumbnail_status"] == "pending"


def test_deleting_a_project_also_deletes_its_thumbnails(db_session_factory):
    """썸네일은 DB가 아니라 디스크에 있어서 cascade로 안 지워진다."""
    project_id = _project(db_session_factory)
    _write_thumbnail(project_id, 1)

    assert client.delete(f"/api/projects/{project_id}").status_code == 200

    assert thumbnail_generator.available_slide_numbers(project_id) == []


# ------------------------------------------------------------ 업로드 연동

def test_upload_queues_thumbnails_and_reports_status(monkeypatch, db_session_factory):
    """업로드 응답은 변환을 기다리지 않고 pending으로 즉시 돌아와야 한다."""
    submitted = {}

    monkeypatch.setattr(thumbnail_generator, "is_available", lambda: True)
    monkeypatch.setattr(main.thumbnail_executor, "submit",
                        lambda fn, *args: submitted.setdefault("args", args))

    response = client.post(
        "/api/projects",
        files={"file": ("deck.pptx", _pptx_bytes(["첫 장", "둘째 장"]),
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )

    try:
        assert response.status_code == 200
        assert response.json()["data"]["thumbnail_status"] == "pending"

        project_id, source_copy = submitted["args"]
        assert project_id == response.json()["project_id"]
        # 백그라운드에 넘길 사본이 실제로 만들어져 있어야 한다. 원본 임시 파일은 요청이
        # 끝나면서 지워지므로, 사본이 없으면 변환할 게 없는 채로 작업이 돈다.
        assert os.path.exists(source_copy)
    finally:
        # submit을 가로챘으니 사본을 지울 백그라운드 작업도 안 돈다. 여기서 치운다.
        for path in submitted.get("args", ())[1:]:
            if isinstance(path, str) and os.path.exists(path):
                os.remove(path)


def test_upload_still_succeeds_when_tools_are_missing(monkeypatch, db_session_factory):
    """⚠️ 핵심. 썸네일 때문에 업로드가 실패하면 안 된다."""
    monkeypatch.setattr(thumbnail_generator, "is_available", lambda: False)

    response = client.post(
        "/api/projects",
        files={"file": ("deck.pptx", _pptx_bytes(["첫 장"]),
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )

    assert response.status_code == 200
    assert response.json()["data"]["thumbnail_status"] == "unavailable"


def test_script_only_project_is_marked_skipped(db_session_factory):
    """대본만 올리면 그릴 원본이 없다. pending으로 두면 영영 안 끝난 것처럼 보인다."""
    response = client.post("/api/projects", data={"script_text": "안녕하세요. 대본만 올립니다."})
    assert response.status_code == 200

    detail = client.get(f"/api/projects/{response.json()['project_id']}").json()["data"]
    assert detail["thumbnail_status"] == "skipped"


def test_background_job_records_status_and_removes_the_copy(monkeypatch, db_session_factory, tmp_path):
    """백그라운드 작업이 상태를 DB에 남기고 임시 사본을 지우는지."""
    project_id = _project(db_session_factory, thumbnail_status="pending")
    source = tmp_path / "copy.pptx"
    source.write_bytes(b"fake")

    monkeypatch.setattr(thumbnail_generator, "generate_for_project",
                        lambda pid, path: {"status": "ready", "count": 3, "message": ""})
    monkeypatch.setattr(main, "job_session_factory", db_session_factory)

    main._run_thumbnail_job(project_id, str(source))

    assert not source.exists(), "임시 사본이 남으면 디스크가 계속 찬다"
    assert client.get(f"/api/projects/{project_id}").json()["data"]["thumbnail_status"] == "ready"


def test_background_job_records_failure(monkeypatch, db_session_factory, tmp_path):
    project_id = _project(db_session_factory, thumbnail_status="pending")
    source = tmp_path / "copy.pptx"
    source.write_bytes(b"fake")

    def _boom(pid, path):
        raise RuntimeError("변환 폭발")

    monkeypatch.setattr(thumbnail_generator, "generate_for_project", _boom)
    monkeypatch.setattr(main, "job_session_factory", db_session_factory)

    main._run_thumbnail_job(project_id, str(source))  # 예외가 새어나오면 안 된다

    assert client.get(f"/api/projects/{project_id}").json()["data"]["thumbnail_status"] == "failed"
    assert not source.exists()
