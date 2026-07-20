import os
import sys
import pypdf
from pptx import Presentation
from pptx.util import Inches, Pt

def convert(pdf_path, pptx_path):
    reader = pypdf.PdfReader(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for i, page in enumerate(reader.pages):
        text = (page.extract_text() or "").strip()
        slide = prs.slides.add_slide(blank_layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.text = f"슬라이드 {i + 1}"
        title_tf.paragraphs[0].font.size = Pt(24)
        title_tf.paragraphs[0].font.bold = True

        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.7))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        if text:
            body_tf.text = text
            for para in body_tf.paragraphs:
                para.font.size = Pt(14)
        else:
            body_tf.text = "(이 페이지는 이미지/캡처라 텍스트를 추출하지 못했습니다)"
            body_tf.paragraphs[0].font.size = Pt(12)
            body_tf.paragraphs[0].font.italic = True

    prs.save(pptx_path)
    print(f"저장 완료: {pptx_path} ({len(reader.pages)} 슬라이드)")

if __name__ == "__main__":
    convert(sys.argv[1], sys.argv[2])
