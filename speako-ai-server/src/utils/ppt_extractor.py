import hashlib
import os
from collections import Counter
from concurrent.futures import ThreadPoolExecutor

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("⚠️ python-pptx 라이브러리가 설치되지 않았습니다. 터미널에서 'pip install python-pptx'를 실행해주세요.")

from claude_vision.claude_ocr_client import ClaudeOcrClient
from clova.vision.image_text_extractor import ImageTextExtractor
from utils.text_heuristics import extract_frequent_terms

# 아이콘/구분선 같은 장식용 이미지는 OCR 대상에서 제외하기 위한 최소 크기 기준(px)
_MIN_OCR_IMAGE_WIDTH = 150
_MIN_OCR_IMAGE_HEIGHT = 100
# HCX-005 비전 입력 제약(가로세로 비율 1:5~5:1)을 벗어나는 길쭉한 이미지도 제외
_MAX_OCR_ASPECT_RATIO = 5.0
# 비전 호출은 슬라이드당 유료 API 호출이다. 한 장에 이미지가 수십 개 박힌 PPT도 있어서
# (실측: 23장짜리에 후보 이미지 97개) 큰 것부터 이만큼만 읽는다.
#
# ⚠️ 3에서 줄이면 안 된다(실측 2026-08-16, 이미지만 있는 14장짜리로 확인).
#       장당 3장 → 슬라이드 10장 / 1,425자
#       장당 2장 → 슬라이드  9장 / 1,207자   (한 장이 통째로 비었다)
#       장당 1장 → 슬라이드  1장 /     6자   (사실상 전멸)
#    **가장 큰 이미지는 대개 배경**이라 글자가 없다. 본문 텍스트는 두세 번째 이미지에
#    들어 있어서, 장수를 줄이면 속도는 빨라져도 대본 근거가 통째로 사라진다.
_MAX_OCR_IMAGES_PER_SLIDE = 3
# 이미 텍스트박스로 내용이 충분히 들어있는 슬라이드는 이미지를 읽어도 얻을 게 별로 없다.
# 비전은 "텍스트가 거의 없는 슬라이드"를 살리는 용도로만 쓴다.
_OCR_SKIP_TEXT_LENGTH = 50
# 비전 호출을 몇 개까지 동시에 보낼지.
#
# 왜 병렬인가(실측 2026-08-16): 글자가 없는 PPT를 한 장씩 순서대로 읽으니 **업로드에만
# 116초**가 걸렸다. 정작 대본 생성은 12초다 — 기다리는 시간의 대부분이 여기였다.
#
# ⚠️ 왜 기본값이 4가 아니라 2인가. 이 구간은 **동시성이 아니라 레이트리밋에 묶여 있다.**
#    같은 파일(이미지 38장)로 잰 값:
#        순차   116초
#        동시4   95.7초   ← 429가 계속 나고 재시도 대기가 23~30초씩 붙는다
#        동시2   86.3초   ← 429가 덜 나서 오히려 빠르다
#    올릴수록 빨라지는 구간이 아니다. 게다가 대본 생성과 같은 HCX 할당량을 나눠 쓰므로,
#    여기서 429를 많이 만들면 대본 생성까지 느려진다. 조절이 필요하면 환경변수로 연다.
_VISION_CONCURRENCY = max(1, int(os.getenv("VISION_MAX_CONCURRENCY", "2")))
# 같은 이미지가 이 수 이상의 슬라이드에 반복되면 템플릿 장식(배경·배너·로고)으로 본다.
#
# 왜 필요한가(실측 2026-08-18, 이미지형 14장): 모든 장에 깔린 1761×773 배경과 963×132
# 배너가 **넓이 상위 3자리를 전부 차지**해서, 정작 내용이 든 이미지(캐릭터 말풍선 등)는
# 비전이 읽지도 못했다. 그 결과 4개 장이 원문 0자가 됐다. 배경은 글자가 없으니 읽어봐야
# 돈만 나가고, 내용 이미지의 자리를 뺏는다.
_TEMPLATE_IMAGE_MIN_SLIDES = 3

class PptExtractor:
    def __init__(self):
        self.image_text_extractor = ImageTextExtractor()
        self.claude_ocr = ClaudeOcrClient()

    def extract_structured_data(self, file_path: str, topic_hint: str = "", outline_hint: str = "") -> dict:
        """
        [업데이트] PPTX 파일 경로를 입력받아 아래와 같이 구조화된 딕셔너리를 반환합니다.
        1. 발표 주제 및 목차/키워드 자동 추출
        2. 슬라이드 번호별 텍스트 완벽 분리
        3. 텍스트박스가 아니라 이미지(캡처/스캔)로만 된 슬라이드는 HCX-005 비전으로 텍스트 추출 시도

        topic_hint / outline_hint: 사용자가 직접 입력한 발표 주제/목차(Figma 유저 플로우 상의 입력값).
        주어지면 이미지 텍스트 인식 시 문맥으로 함께 제공해 정확도를 높인다. 없으면 힌트 없이 읽는다.
        """
        if not os.path.exists(file_path):
            print(f"❌ 파일을 찾을 수 없습니다: {file_path}")
            return {"metadata": {"topic": "", "keywords": []}, "slides": []}

        context_hint = ""
        if topic_hint:
            context_hint += f"발표 주제: {topic_hint}\n"
        if outline_hint:
            context_hint += f"목차/가이드라인: {outline_hint}\n"

        try:
            prs = Presentation(file_path)
            slides_data = []
            front_text_for_analysis = []
            all_slide_texts = []

            # ── 1단계: 파일을 훑어 텍스트를 모으고, 비전으로 읽을 이미지 후보를 쌓는다.
            #    python-pptx 객체를 워커 스레드로 넘기지 않으려고 여기서 blob까지 꺼내둔다.
            raw = []            # 슬라이드별 {"texts": [...], "candidates": [(넓이, blob, type), ...]}
            blob_slide_count = Counter()   # 이미지 해시 → 등장한 슬라이드 수 (템플릿 감지용)
            for slide in prs.slides:
                slide_texts = []
                ocr_candidates = []
                seen_in_slide = set()
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        # 아이콘/구분선 같은 장식용 이미지는 크기/비율로 걸러내고,
                        # 내용이 있을 법한 이미지만 후보로 남긴다.
                        try:
                            image = shape.image
                            width_px, height_px = image.size
                            blob = image.blob
                        except Exception:
                            continue
                        if width_px < _MIN_OCR_IMAGE_WIDTH or height_px < _MIN_OCR_IMAGE_HEIGHT:
                            continue
                        aspect_ratio = max(width_px, height_px) / max(1, min(width_px, height_px))
                        if aspect_ratio > _MAX_OCR_ASPECT_RATIO:
                            continue
                        digest = hashlib.sha1(blob).hexdigest()
                        # 같은 슬라이드에 같은 이미지가 두 번 있어도 슬라이드 수는 1로 센다.
                        if digest not in seen_in_slide:
                            seen_in_slide.add(digest)
                            blob_slide_count[digest] += 1
                        ocr_candidates.append(
                            (width_px * height_px, digest, blob, image.content_type))

                raw.append({"texts": slide_texts, "candidates": ocr_candidates})

            # ── 1.5단계: 여러 장에 반복되는 이미지는 템플릿 장식이다 — OCR 후보에서 뺀다.
            #    배경·배너가 넓이 상위 자리를 차지하면 내용 이미지가 읽히지 못한다(위 상수 주석).
            template_digests = {
                digest for digest, count in blob_slide_count.items()
                if count >= _TEMPLATE_IMAGE_MIN_SLIDES
            }

            pending = []   # 슬라이드별 {"texts": [...], "images": [(blob, content_type), ...]}
            for entry in raw:
                slide_texts = entry["texts"]
                candidates = [c for c in entry["candidates"] if c[1] not in template_digests]
                # 전부 템플릿뿐이면(표지가 배경만으로 된 장 등) 원래 후보로 되돌린다 —
                # 아예 안 읽는 것보다는 배경이라도 읽어보는 쪽이 낫다.
                if not candidates and entry["candidates"]:
                    candidates = entry["candidates"]

                # 텍스트가 거의 없는 슬라이드(= 이미지가 곧 내용인 장표)만 비전으로 읽는다.
                # 큰 이미지일수록 본문일 확률이 높으니 넓이 순으로 상위 몇 장만 본다.
                images, extra_images = [], []
                if candidates and len("".join(slide_texts)) < _OCR_SKIP_TEXT_LENGTH:
                    candidates.sort(key=lambda item: item[0], reverse=True)
                    ordered = [(blob, content_type) for _, _, blob, content_type in candidates]
                    images = ordered[:_MAX_OCR_IMAGES_PER_SLIDE]
                    # 상위 3장에서 글자가 하나도 안 나오면 읽을 예비 후보 (아래 2차 라운드).
                    extra_images = ordered[_MAX_OCR_IMAGES_PER_SLIDE:_MAX_OCR_IMAGES_PER_SLIDE * 2]

                pending.append({"texts": slide_texts, "images": images, "extra_images": extra_images})

            # ── 2단계: 모아둔 이미지를 한꺼번에 동시 호출한다.
            #    슬라이드를 하나씩 기다리면 이미지 30장짜리에서 2분이 걸린다(위 상수 주석 참고).
            def _read(job):
                s_idx, i_idx, blob, content_type = job
                try:
                    return (s_idx, i_idx), self.image_text_extractor.extract_text_from_image(
                        blob, context_hint, content_type
                    )
                except Exception as err:
                    # 이미지 한 장 실패가 업로드 전체를 깨면 안 된다. 그 장만 비우고 간다.
                    print(f"⚠️ 이미지 텍스트 인식 실패(슬라이드 {s_idx + 1}): {err}")
                    return (s_idx, i_idx), ""

            def _run_jobs(jobs, round_label):
                if not jobs:
                    return
                print(f"🖼️  이미지 {len(jobs)}장을 동시 {_VISION_CONCURRENCY}개씩 읽습니다.{round_label}")
                with ThreadPoolExecutor(max_workers=_VISION_CONCURRENCY) as pool:
                    for key, text in pool.map(_read, jobs):
                        if text and text.strip():
                            ocr_results[key] = text.strip()

            ocr_results = {}
            _run_jobs([(s_idx, i_idx, blob, content_type)
                       for s_idx, entry in enumerate(pending)
                       for i_idx, (blob, content_type) in enumerate(entry["images"])], "")

            # ── 2차 라운드: 상위 3장에서 글자가 하나도 안 나온 장만 예비 후보를 마저 읽는다.
            #
            # 왜 필요한가(실측 2026-08-19, 이미지형 2·3·13장): 텍스트가 든 이미지(캐릭터 카드의
            # 말풍선 "반갑습니다 진순입니다:)")가 넓이 4~5번째라 1차에서 읽히지 않았고, 그 장들은
            # 원문 0자가 됐다. 모든 장의 예비까지 다 읽으면 비용이 배가 되지만, **빈 장만**
            # 추가로 읽으면 그 장들에서만 호출이 늘어난다.
            retry_jobs = []
            for s_idx, entry in enumerate(pending):
                if entry["texts"] or not entry["extra_images"]:
                    continue
                if any(ocr_results.get((s_idx, j)) for j in range(len(entry["images"]))):
                    continue
                base = len(entry["images"])
                retry_jobs.extend(
                    (s_idx, base + j, blob, content_type)
                    for j, (blob, content_type) in enumerate(entry["extra_images"]))
            _run_jobs(retry_jobs, " (빈 장 2차)")
            # 병합 루프가 예비 결과까지 훑도록 이미지 목록을 합쳐둔다.
            for entry in pending:
                entry["images"] = entry["images"] + entry["extra_images"]

            # ── 2.5차: HCX가 한 글자도 못 읽은 장은 Claude OCR로 다시 읽어본다 (키 있을 때만).
            #
            # HCX-005 비전은 큰 일반 폰트도 못 읽는 경우가 있는데(2026-08-19 실험) Claude는
            # 같은 이미지를 완벽히 읽었다. 대본 생성은 그대로 HCX다 — 여기는 글자 전사만.
            # 키(ANTHROPIC_API_KEY)가 없으면 claude_ocr가 조용히 빈 값을 돌려줘 3차로 넘어간다.
            if not self.claude_ocr.use_fallback:
                for s_idx, entry in enumerate(pending):
                    if entry["texts"] or not entry["images"]:
                        continue
                    if any(ocr_results.get((s_idx, j)) for j in range(len(entry["images"]))):
                        continue
                    for j, (blob, content_type) in enumerate(entry["images"]):
                        text = self.claude_ocr.extract_text_from_image(blob, content_type)
                        if text:
                            ocr_results[(s_idx, j)] = text
                    if any(ocr_results.get((s_idx, j)) for j in range(len(entry["images"]))):
                        print(f"🔁 슬라이드 {s_idx + 1}: HCX가 못 읽은 글자를 Claude OCR이 읽었습니다.")

            # ── 3차: 그래도 글자가 없는 장은 "어떤 장면인지" 한 문장을 받아 [화면 묘사]로 저장.
            #
            # HCX-005 비전은 큰 일반 폰트도 못 읽는 경우가 있다(2026-08-19 실험 — nextStep
            # 백로그 참고). 글자를 포기한 장이라도 장면(인물이 인사하는 자기소개 장 등)을 알면
            # 생성기가 장의 역할에 맞는 대본을 쓸 수 있다. 가장 큰 이미지 한 장만 물어본다.
            for s_idx, entry in enumerate(pending):
                if entry["texts"] or not entry["images"]:
                    continue
                if any(ocr_results.get((s_idx, j)) for j in range(len(entry["images"]))):
                    continue
                # 가장 큰 이미지가 빈 말풍선 같은 무정보 그림일 수 있어(실측: 이미지형 3장),
                # 쓸 만한 묘사가 나올 때까지 상위 3장까지 시도한다.
                for blob, content_type in entry["images"][:3]:
                    try:
                        scene = self.image_text_extractor.describe_scene(blob, content_type)
                    except Exception as err:
                        print(f"⚠️ 장면 설명 실패(슬라이드 {s_idx + 1}): {err}")
                        scene = ""
                    if scene:
                        entry["scene"] = scene
                        break

            # ── 3단계: 원래 순서대로 병합한다(텍스트박스 먼저, 그 뒤에 넓이 순 이미지).
            for i, entry in enumerate(pending):
                slide_texts = list(entry["texts"])
                for i_idx in range(len(entry["images"])):
                    text = ocr_results.get((i, i_idx))
                    if text:
                        slide_texts.append(text)

                # 글자가 없고 장면 설명만 있으면 라벨을 붙여 격리한다. 라벨 없이 넣으면
                # 생성기가 묘사를 슬라이드 원문으로 믿는다(과거 환각 사고의 원인).
                if not slide_texts and entry.get("scene"):
                    slide_texts = [f"[화면 묘사] {entry['scene']}"]

                slide_content = "\n".join(slide_texts)

                # 내용이 안 뽑힌 장도 **버리지 않고 빈 내용으로 유지한다.**
                # 예전엔 여기서 걸러냈는데, 그러면 원본 14장이 10장으로 줄어 사용자는
                # "내 PPT에서 장이 사라졌다"를 보게 된다(실측: 02_이미지형.pptx 14장 → 10장).
                # 장 번호가 원본과 밀리면 썸네일·하이라이트 좌표까지 연쇄로 어긋난다.
                # 빈 장의 대본은 생성기의 근거 없는 장 처리(_is_thin_source)가 짧게 맡는다.
                slides_data.append({
                    "slide_number": i + 1,
                    "content": slide_content
                })

                if slide_content.strip():
                    all_slide_texts.extend(slide_texts)

                    # 주제 및 목차 탐지를 위해 초반 1~3장 텍스트는 별도로도 수집
                    if i < 3:
                        front_text_for_analysis.extend(slide_texts)

            # [요구사항 4 반영] 초반 슬라이드 텍스트 기반 주제 추출 + 전체 슬라이드 기반 키워드 추출
            topic, keywords = self._extract_metadata(front_text_for_analysis, all_slide_texts)

            print(f"✅ PPT 구조화 추출 완료! (총 {len(slides_data)}장 분석)")
            return {
                "metadata": {
                    "topic": topic,
                    "keywords": keywords
                },
                "slides": slides_data
            }

        except Exception as e:
            print(f"❌ PPT 추출 중 오류 발생: {e}")
            return {"metadata": {"topic": "", "keywords": []}, "slides": []}

    def _extract_metadata(self, front_texts: list, all_texts: list = None):
        """초반 슬라이드 텍스트로 발표 주제를, 전체 슬라이드 텍스트로 키워드를 휴리스틱하게 추론합니다."""
        if not front_texts:
            return "주제 미상", []

        all_texts = all_texts if all_texts else front_texts

        # 보통 첫 번째 텍스트 덩어리가 발표 주제(Title)일 확률이 높음
        topic = front_texts[0].strip()
        keywords = []

        # '목차', 'index', 'agenda' 등의 단어가 포함된 슬라이드가 있다면 그 직후 텍스트를 키워드로 간주
        for i, text in enumerate(all_texts):
            lower_text = text.lower()
            if any(keyword in lower_text for keyword in ['목차', 'index', 'contents', '순서', 'agenda', '차례', 'outline']):
                keywords = all_texts[i + 1: i + 6]
                break

        # 목차 슬라이드가 없다면, 전체 슬라이드 텍스트에서 자주 등장하는 단어를 키워드로 추론
        if not keywords:
            keywords = self._extract_keywords_by_frequency(all_texts, topic)

        return topic, keywords

    def _extract_keywords_by_frequency(self, texts: list, topic: str, top_n: int = 5) -> list:
        """목차 슬라이드가 없는 PPT를 위한 대체 키워드 추출: 전체 텍스트에서 자주 등장하는 단어를 뽑는다."""
        return extract_frequent_terms(texts, exclude={topic}, top_n=top_n)

# ==========================================
# 🧪 [테스트 코드]
# ==========================================
if __name__ == "__main__":
    extractor = PptExtractor()
    test_file = "sample.pptx" 
    
    if os.path.exists(test_file):
        # 기존 통글 추출 대신 구조화된 데이터 추출 메서드 사용
        result_data = extractor.extract_structured_data(test_file)
        
        import json
        print("\n✨ [추출된 PPT 구조화 데이터] ✨")
        print(json.dumps(result_data, indent=2, ensure_ascii=False))
    else:
        print(f"⚠️ '{test_file}' 파일이 존재하지 않아 추출 테스트를 건너뜁니다.")